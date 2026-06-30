from __future__ import annotations

from pathlib import Path
from typing import Any, Dict, List, Tuple

import cv2
import numpy as np
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR
from pptx.util import Inches, Pt

from .element_extractor import ImageElement, LineElement, SlideElements
from .ocr_engine import TextElement


def px_to_in(x: float, total_px: int, total_in: float) -> float:
    return float(x) / float(total_px) * float(total_in)


def rgb_to_color(rgb: Tuple[int, int, int]) -> RGBColor:
    r, g, b = rgb
    return RGBColor(max(0, min(255, int(r))), max(0, min(255, int(g))), max(0, min(255, int(b))))


def create_blank_prs(cfg: Dict[str, Any]) -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(float(cfg.get("slide_width_in", 13.333)))
    prs.slide_height = Inches(float(cfg.get("slide_height_in", 7.5)))
    return prs


def add_reference_image_slide(prs: Presentation, image_path: str | Path, cfg: Dict[str, Any]) -> None:
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    slide.shapes.add_picture(str(image_path), 0, 0, width=prs.slide_width, height=prs.slide_height)


def build_reference_ppt(image_paths: List[str | Path], out_path: str | Path, cfg: Dict[str, Any]) -> None:
    prs = create_blank_prs(cfg)
    for p in image_paths:
        add_reference_image_slide(prs, p, cfg)
    prs.save(str(out_path))


def _font_size_from_bbox(t: TextElement, cfg: Dict[str, Any]) -> float:
    canvas_h = int(cfg.get("canvas_height", 1080))
    slide_h = float(cfg.get("slide_height_in", 7.5))
    # OCR 框高度 → 页面英寸 → point。中文行高通常大于字号，所以取 0.62～0.75。
    factor = float(cfg.get("rebuild", {}).get("font_size_factor", 0.68))
    pt = t.h / canvas_h * slide_h * 72.0 * factor
    return max(float(cfg.get("rebuild", {}).get("min_font_pt", 6.0)), min(float(cfg.get("rebuild", {}).get("max_font_pt", 40.0)), pt))


def _clip(v: int, lo: int, hi: int) -> int:
    return max(lo, min(hi, int(v)))


def _roi_background_rgb(bg_bgr: np.ndarray | None, t: TextElement, pad: int = 2) -> Tuple[int, int, int]:
    """
    估计“可编辑文字框正下方”的局部背景色，返回 RGB。
    注意：这里不要用 OCR 识别到的原文字颜色，而是要看文字将要落在什么背景上。
    """
    if bg_bgr is None or bg_bgr.size == 0:
        return (255, 255, 255)
    H, W = bg_bgr.shape[:2]
    x0 = _clip(t.x - pad, 0, W - 1)
    y0 = _clip(t.y - pad, 0, H - 1)
    x1 = _clip(t.x + t.w + pad, 0, W)
    y1 = _clip(t.y + t.h + pad, 0, H)
    if x1 <= x0 or y1 <= y0:
        return (255, 255, 255)
    roi = bg_bgr[y0:y1, x0:x1]
    if roi.size == 0:
        return (255, 255, 255)

    # 使用中位数而不是均值，避免残留文字笔画、噪点、边缘线影响背景估计。
    bgr = np.median(roi.reshape(-1, 3), axis=0)
    return (int(bgr[2]), int(bgr[1]), int(bgr[0]))


def _relative_luminance(rgb: Tuple[int, int, int]) -> float:
    r, g, b = [x / 255.0 for x in rgb]
    # sRGB relative luminance 的简化 gamma 前版本，在这里用于黑/白阈值足够。
    return 0.2126 * r + 0.7152 * g + 0.0722 * b


def _saturation(rgb: Tuple[int, int, int]) -> float:
    r, g, b = [x / 255.0 for x in rgb]
    mx, mn = max(r, g, b), min(r, g, b)
    return 0.0 if mx <= 1e-6 else (mx - mn) / mx


def _contrast_ratio(rgb1: Tuple[int, int, int], rgb2: Tuple[int, int, int]) -> float:
    l1 = _relative_luminance(rgb1)
    l2 = _relative_luminance(rgb2)
    hi, lo = max(l1, l2), min(l1, l2)
    return (hi + 0.05) / (lo + 0.05)


def _choose_contrast_text_color(bg_rgb: Tuple[int, int, int], cfg: Dict[str, Any]) -> Tuple[int, int, int]:
    """
    根据局部背景色选择可编辑文字颜色。

    本版按你的要求做：不要追求还原原字体颜色，而是追求“在背景上足够突出”。
    默认策略是黑/白二选一：
    - 背景是绿色、蓝色、紫色、红色、深色、饱和色 → 用白字；
    - 背景是白色、浅灰、浅黄、很亮区域 → 用黑字。

    也支持 true_inverse，但默认不用。因为绿色的数学补色是洋红，不适合 PPT 阅读；
    你举的例子里绿色背景、蓝色背景都希望用白字，所以这里默认是 high_contrast_bw。
    """
    rebuild = cfg.get("rebuild", {})
    mode = str(rebuild.get("contrast_text_mode", "high_contrast_bw")).lower()
    if mode in ("inverse", "true_inverse", "rgb_inverse"):
        return tuple(255 - int(c) for c in bg_rgb)  # type: ignore[return-value]

    white = (255, 255, 255)
    black = (0, 0, 0)

    # 先用标准对比度判断黑/白哪个更明显。
    cr_white = _contrast_ratio(white, bg_rgb)
    cr_black = _contrast_ratio(black, bg_rgb)

    lum = _relative_luminance(bg_rgb) * 255.0
    sat = _saturation(bg_rgb)

    # 你明确提到绿色、蓝色背景要白字：很多投影色块的亮度并不低，
    # 单纯按 luminance 可能会选黑字，所以额外加入“饱和色块优先白字”。
    color_bg_white_luma_limit = float(rebuild.get("color_bg_white_luma_limit", 220.0))
    color_bg_saturation_threshold = float(rebuild.get("color_bg_saturation_threshold", 0.18))
    dark_bg_luma_threshold = float(rebuild.get("dark_bg_luma_threshold", 170.0))

    if lum < dark_bg_luma_threshold:
        return white
    if sat >= color_bg_saturation_threshold and lum < color_bg_white_luma_limit:
        return white

    # 对接近白色、浅灰色、浅黄等背景，黑字通常最清楚。
    return black if cr_black >= cr_white else white


def _text_color_for_box(t: TextElement, cfg: Dict[str, Any], bg_bgr: np.ndarray | None) -> Tuple[int, int, int]:
    rebuild = cfg.get("rebuild", {})
    policy = str(rebuild.get("text_color_policy", "contrast_with_background")).lower()
    if policy in ("ocr", "original", "estimated_original"):
        return t.color
    bg_rgb = _roi_background_rgb(bg_bgr, t, int(rebuild.get("background_sample_padding_px", 3)))
    return _choose_contrast_text_color(bg_rgb, cfg)


def _add_text(slide, t: TextElement, prs: Presentation, cfg: Dict[str, Any], bg_bgr: np.ndarray | None = None) -> None:
    cw = int(cfg.get("canvas_width", 1920))
    ch = int(cfg.get("canvas_height", 1080))
    sw = float(cfg.get("slide_width_in", 13.333))
    sh = float(cfg.get("slide_height_in", 7.5))
    rebuild = cfg.get("rebuild", {})
    pad_px = float(rebuild.get("textbox_padding_px", 4))

    left = Inches(px_to_in(max(0, t.x - pad_px), cw, sw))
    top = Inches(px_to_in(max(0, t.y - pad_px), ch, sh))
    width = Inches(px_to_in(max(t.w + pad_px * 2, 12), cw, sw))
    height = Inches(px_to_in(max(t.h + pad_px * 2.4, 14), ch, sh))

    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.word_wrap = False
    tf.vertical_anchor = MSO_ANCHOR.TOP

    p = tf.paragraphs[0]
    p.space_before = Pt(0)
    p.space_after = Pt(0)
    run = p.add_run()
    run.text = t.text
    run.font.name = str(rebuild.get("font_name", "Microsoft YaHei"))
    run.font.size = Pt(_font_size_from_bbox(t, cfg))
    run.font.color.rgb = rgb_to_color(_text_color_for_box(t, cfg, bg_bgr))


def _add_image(slide, elem: ImageElement, prs: Presentation, cfg: Dict[str, Any]) -> None:
    cw = int(cfg.get("canvas_width", 1920))
    ch = int(cfg.get("canvas_height", 1080))
    sw = float(cfg.get("slide_width_in", 13.333))
    sh = float(cfg.get("slide_height_in", 7.5))
    slide.shapes.add_picture(
        elem.path,
        Inches(px_to_in(elem.x, cw, sw)),
        Inches(px_to_in(elem.y, ch, sh)),
        width=Inches(px_to_in(elem.w, cw, sw)),
        height=Inches(px_to_in(elem.h, ch, sh)),
    )


def _add_line(slide, elem: LineElement, prs: Presentation, cfg: Dict[str, Any]) -> None:
    cw = int(cfg.get("canvas_width", 1920))
    ch = int(cfg.get("canvas_height", 1080))
    sw = float(cfg.get("slide_width_in", 13.333))
    sh = float(cfg.get("slide_height_in", 7.5))
    x1 = Inches(px_to_in(elem.x1, cw, sw))
    y1 = Inches(px_to_in(elem.y1, ch, sh))
    x2 = Inches(px_to_in(elem.x2, cw, sw))
    y2 = Inches(px_to_in(elem.y2, ch, sh))
    line = slide.shapes.add_connector(1, x1, y1, x2, y2)
    line.line.color.rgb = rgb_to_color(elem.color)
    line.line.width = Pt(elem.width)


def add_rebuilt_slide(prs: Presentation, elems: SlideElements, cfg: Dict[str, Any]) -> None:
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    bg_bgr = None
    if elems.background_image:
        slide.shapes.add_picture(str(elems.background_image), 0, 0, width=prs.slide_width, height=prs.slide_height)
        bg_bgr = cv2.imread(str(elems.background_image))
    else:
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = rgb_to_color(elems.bg_color)

    for img in elems.images:
        _add_image(slide, img, prs, cfg)
    for line in elems.lines:
        _add_line(slide, line, prs, cfg)
    for text in elems.texts:
        _add_text(slide, text, prs, cfg, bg_bgr=bg_bgr)


def build_editable_ppt(slides: List[SlideElements], out_path: str | Path, cfg: Dict[str, Any]) -> None:
    prs = create_blank_prs(cfg)
    for elems in slides:
        add_rebuilt_slide(prs, elems, cfg)
    Path(out_path).parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
